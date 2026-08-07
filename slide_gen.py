"""
Slayd (PowerPoint) tayyorlovchi modul.

Oqim:
  1. Gemini API - mavzu + slaydlar soniga qarab har bir slayd uchun
     sarlavha/matn/rasm kalit so'zini generatsiya qiladi (matn, bepul reja).
  2. Pexels API - har bir slayd uchun kalit so'z asosida mos fotosurat topadi
     (tayyor stok fotosurat, bepul reja).
  3. python-pptx - tanlangan dizayn shabloniga (4 tadan biriga) qarab hammasini
     PPTX faylga yig'adi.

Og'ir mahalliy AI model yo'q - faqat ikkita tashqi API'ga qisqa so'rov,
shuning uchun server uchun yengil.
"""
import io
import json
import logging
import os

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
PEXELS_API_KEY = os.getenv("PEXELS_API_KEY", "")

GEMINI_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-flash-latest:generateContent"
PEXELS_SEARCH_URL = "https://api.pexels.com/v1/search"

MIN_SLIDES = 5
MAX_SLIDES = 30
PRICE_PER_SLIDE = 500

TEMPLATES = {
    "blue": "Ko'k klassik",
    "dark": "Tungi tex",
    "warm": "Iliq ijodiy",
    "yellow": "Sariq hisobot",
}


# ─────────────────────────────────────────────────────────
# 1-qadam: Gemini - slaydlar tarkibini generatsiya qilish
# ─────────────────────────────────────────────────────────

def generate_outline(topic: str, num_slides: int) -> list:
    """Gemini API orqali mavzuga mos slaydlar tarkibini (JSON) generatsiya qiladi."""
    if not GEMINI_API_KEY:
        raise RuntimeError("GEMINI_API_KEY sozlanmagan. Render'da Environment bo'limiga qo'shing.")

    prompt = f"""Sen professional taqdimot (prezentatsiya) tayyorlovchi yordamchisan.
Mavzu: "{topic}"
Jami slaydlar soni: {num_slides}

FAQAT quyidagi JSON massivini qaytar (boshqa hech qanday matn, izoh yoki ``` belgisi yozma):

[
  {{"type": "title", "title": "...", "subtitle": "..."}},
  {{"type": "content", "title": "...", "bullets": ["...", "...", "..."], "image_keyword": "english keyword"}},
  {{"type": "divider", "title": "Xulosa", "subtitle": "..."}}
]

Qat'iy qoidalar:
- Jami roppa-rosa {num_slides} ta element (slayd) bo'lsin
- 1-slayd har doim "title" turida (mavzu sarlavhasi + qisqa subtitle)
- Oxirgi slayd "divider" turida ("Xulosa" yoki "Yakun" kabi)
- Qolgan barcha slaydlar "content" turida: title (qisqa), 5-6 ta bullets (har biri 10-18 so'z, mazmunli va aniq faktlar/tushuntirishlar bilan), va image_keyword
- Matn (title, subtitle, bullets) - o'zbek tilida
- image_keyword - FAQAT ingliz tilida, 1-3 so'z, Pexels fotosurat qidiruvi uchun mos (masalan "business meeting", "mountain landscape")
- Bullets qisqa, aniq, taqdimotga mos uslubda bo'lsin, gap oxirida nuqta qo'yma
"""

    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {"temperature": 0.7, "maxOutputTokens": 32768},
    }

    resp = requests.post(f"{GEMINI_URL}?key={GEMINI_API_KEY}", json=payload, timeout=90)
    if resp.status_code != 200:
        raise RuntimeError(f"Gemini xatosi ({resp.status_code}): {resp.text[:300]}")

    data = resp.json()
    try:
        text = data["candidates"][0]["content"]["parts"][0]["text"].strip()
    except (KeyError, IndexError) as e:
        raise RuntimeError(f"Gemini javobini o'qib bo'lmadi: {data}") from e

    # Model ba'zan ```json bilan o'raydi - tozalaymiz
    if text.startswith("```"):
        text = text.strip("`")
        if text.startswith("json"):
            text = text[4:]

    try:
        slides = json.loads(text.strip())
    except json.JSONDecodeError as e:
        raise RuntimeError(f"Gemini JSON qaytarmadi: {text[:300]}") from e

    if not isinstance(slides, list) or not slides:
        raise RuntimeError("Gemini bo'sh yoki noto'g'ri natija qaytardi")

    return slides


# ─────────────────────────────────────────────────────────
# 2-qadam: Pexels - mavzuga mos fotosurat topish
# ─────────────────────────────────────────────────────────

def fetch_pexels_image(keyword: str) -> bytes:
    """Pexels'dan kalit so'zga mos fotosuratni yuklab oladi. Topilmasa None."""
    if not PEXELS_API_KEY or not keyword:
        return None
    try:
        resp = requests.get(
            PEXELS_SEARCH_URL,
            params={"query": keyword, "per_page": 1},
            headers={"Authorization": PEXELS_API_KEY},
            timeout=20,
        )
        if resp.status_code != 200:
            logger.warning(f"Pexels xatosi ({resp.status_code}) '{keyword}' uchun")
            return None
        photos = resp.json().get("photos", [])
        if not photos:
            return None
        img_url = photos[0]["src"]["large"]
        img_resp = requests.get(img_url, timeout=20)
        img_resp.raise_for_status()
        return img_resp.content
    except Exception:
        logger.exception(f"Pexels'dan rasm olishda xato: {keyword}")
        return None


# ─────────────────────────────────────────────────────────
# 3-qadam: python-pptx - dizayn shablonlari
# ─────────────────────────────────────────────────────────

def _set_no_line(shape):
    shape.line.fill.background()


def _add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Arial"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    return tb


def _add_bullets(slide, x, y, w, h, items, size, color, font="Arial"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = "•  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    return tb


def _bg_fill(prs, slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    _set_no_line(bg)
    return bg


def _add_image_or_fallback(slide, image_bytes, x, y, w, h, fallback_color):
    """Rasm bo'lsa joylaydi (to'ldirib kesib), bo'lmasa rangli to'rtburchak qo'yadi."""
    if image_bytes:
        try:
            slide.shapes.add_picture(io.BytesIO(image_bytes), x, y, width=w, height=h)
            return
        except Exception:
            logger.exception("Rasmni pptx'ga joylashda xato")
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fallback_color
    _set_no_line(box)


# ---------- Shablon 1: Ko'k klassik ----------

def _render_blue(prs, slides):
    BLUE = RGBColor(0x1B, 0x4F, 0x72)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x1A, 0x1A, 0x1A)
    GRAY = RGBColor(0x5A, 0x5A, 0x5A)
    PLACEHOLDER = RGBColor(0xD9, 0xE2, 0xEA)
    blank = prs.slide_layouts[6]

    for i, sd in enumerate(slides):
        s = prs.slides.add_slide(blank)
        _bg_fill(prs, s, WHITE)
        t = sd.get("type", "content")

        if t == "title":
            blk = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-2), Inches(5.8), Inches(9), Inches(5))
            blk.rotation = -8
            blk.fill.solid(); blk.fill.fore_color.rgb = BLUE; _set_no_line(blk)
            _add_text(s, Inches(0.9), Inches(2.3), Inches(9.8), Inches(1.8), sd.get("title", ""), 38, DARK, bold=True)
            _add_text(s, Inches(0.9), Inches(4.3), Inches(9.5), Inches(0.6), sd.get("subtitle", ""), 18, GRAY)
            _add_text(s, Inches(0.9), Inches(7.0), Inches(6), Inches(0.5), "MADESY", 13, WHITE, bold=True)

        elif t == "divider":
            _bg_fill(prs, s, BLUE)
            _add_text(s, Inches(1), Inches(2.7), Inches(11.3), Inches(1.6), sd.get("title", ""), 40, WHITE, bold=True)
            _add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.6), sd.get("subtitle", ""), 18, RGBColor(0xC9, 0xDD, 0xEA))

        else:
            band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.6), prs.slide_height)
            band.fill.solid(); band.fill.fore_color.rgb = BLUE; _set_no_line(band)
            _add_image_or_fallback(s, sd.get("_image"), Inches(0.7), Inches(2.1), Inches(3.2), Inches(3.2), PLACEHOLDER)
            _add_text(s, Inches(5.1), Inches(0.8), Inches(7.4), Inches(1.0), sd.get("title", ""), 28, DARK, bold=True)
            _add_bullets(s, Inches(5.1), Inches(2.0), Inches(7.4), Inches(4.5), sd.get("bullets", []), 16, DARK)


# ---------- Shablon 2: Tungi tex ----------

def _render_dark(prs, slides):
    NAVY = RGBColor(0x12, 0x0E, 0x2B)
    PURPLE = RGBColor(0x4A, 0x2F, 0x7A)
    ACCENT = RGBColor(0x8B, 0x5C, 0xF6)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT = RGBColor(0xC9, 0xB8, 0xE8)
    PLACEHOLDER = RGBColor(0x2A, 0x22, 0x4A)
    blank = prs.slide_layouts[6]

    for sd in slides:
        s = prs.slides.add_slide(blank)
        t = sd.get("type", "content")

        if t == "title":
            _bg_fill(prs, s, NAVY)
            glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5), Inches(5))
            glow.fill.solid(); glow.fill.fore_color.rgb = PURPLE; _set_no_line(glow)
            _add_text(s, Inches(0.9), Inches(2.2), Inches(10.5), Inches(1.8), sd.get("title", ""), 36, WHITE, bold=True)
            _add_text(s, Inches(0.9), Inches(4.2), Inches(9), Inches(0.6), sd.get("subtitle", ""), 18, LIGHT)
            _add_text(s, Inches(0.9), Inches(6.5), Inches(6), Inches(0.5), "MADESY", 13, ACCENT, bold=True)

        elif t == "divider":
            _bg_fill(prs, s, PURPLE)
            _add_text(s, Inches(1), Inches(2.7), Inches(11.3), Inches(1.6), sd.get("title", ""), 40, WHITE, bold=True)
            _add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.6), sd.get("subtitle", ""), 18, LIGHT)

        else:
            _bg_fill(prs, s, NAVY)
            _add_image_or_fallback(s, sd.get("_image"), Inches(8.3), Inches(1.2), Inches(4.1), Inches(5.1), PLACEHOLDER)
            _add_text(s, Inches(0.9), Inches(0.9), Inches(7), Inches(1.0), sd.get("title", ""), 26, WHITE, bold=True)
            _add_bullets(s, Inches(0.9), Inches(2.1), Inches(7), Inches(4.5), sd.get("bullets", []), 16, LIGHT)


# ---------- Shablon 3: Iliq ijodiy ----------

def _render_warm(prs, slides):
    CORAL = RGBColor(0xE8, 0x7A, 0x5D)
    TEAL = RGBColor(0x3D, 0x8B, 0x8A)
    CREAM = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x2B, 0x2B, 0x2B)
    GRAY = RGBColor(0x6B, 0x6B, 0x6B)
    PLACEHOLDER = RGBColor(0xF3, 0xDA, 0xD0)
    blank = prs.slide_layouts[6]

    content_num = 0
    for i, sd in enumerate(slides):
        s = prs.slides.add_slide(blank)
        t = sd.get("type", "content")

        if t == "title":
            _bg_fill(prs, s, CREAM)
            c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(-1.2), Inches(4), Inches(4))
            c1.fill.solid(); c1.fill.fore_color.rgb = CORAL; _set_no_line(c1)
            c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(5.5), Inches(4), Inches(4))
            c2.fill.solid(); c2.fill.fore_color.rgb = TEAL; _set_no_line(c2)
            _add_text(s, Inches(0.9), Inches(2.5), Inches(10.5), Inches(1.8), sd.get("title", ""), 34, DARK, bold=True)
            _add_text(s, Inches(0.9), Inches(4.4), Inches(9), Inches(0.6), sd.get("subtitle", ""), 18, GRAY)

        elif t == "divider":
            _bg_fill(prs, s, CORAL)
            _add_text(s, Inches(1), Inches(2.7), Inches(11.3), Inches(1.6), sd.get("title", ""), 40, CREAM, bold=True)
            _add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.6), sd.get("subtitle", ""), 18, RGBColor(0xFF, 0xE8, 0xDF))

        else:
            content_num += 1
            _bg_fill(prs, s, CREAM)
            badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.7), Inches(1.3), Inches(0.5))
            badge.fill.solid(); badge.fill.fore_color.rgb = TEAL if content_num % 2 else CORAL; _set_no_line(badge)
            _add_text(s, Inches(0.9), Inches(0.82), Inches(1.3), Inches(0.4), f"{content_num:02d}", 15, CREAM, bold=True, align=PP_ALIGN.CENTER)
            _add_image_or_fallback(s, sd.get("_image"), Inches(8.6), Inches(1.3), Inches(3.8), Inches(3.8), PLACEHOLDER)
            _add_text(s, Inches(0.9), Inches(1.55), Inches(7), Inches(1.0), sd.get("title", ""), 26, DARK, bold=True)
            _add_bullets(s, Inches(0.9), Inches(2.7), Inches(7.3), Inches(4.0), sd.get("bullets", []), 15, DARK)


# ---------- Shablon 4: Sariq hisobot ----------

def _render_yellow(prs, slides):
    YELLOW = RGBColor(0xF5, 0xB7, 0x1E)
    BLACK = RGBColor(0x1A, 0x1A, 0x1A)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY_LIGHT = RGBColor(0xE5, 0xE5, 0xE5)
    PLACEHOLDER = RGBColor(0x3A, 0x3A, 0x3A)
    blank = prs.slide_layouts[6]

    for sd in slides:
        s = prs.slides.add_slide(blank)
        t = sd.get("type", "content")

        if t == "title":
            _bg_fill(prs, s, BLACK)
            blk = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.6), prs.slide_height)
            blk.fill.solid(); blk.fill.fore_color.rgb = YELLOW; _set_no_line(blk)
            _add_text(s, Inches(5.2), Inches(2.4), Inches(7.3), Inches(1.8), sd.get("title", ""), 36, WHITE, bold=True)
            _add_text(s, Inches(5.2), Inches(4.4), Inches(7), Inches(0.6), sd.get("subtitle", ""), 18, GRAY_LIGHT)

        elif t == "divider":
            _bg_fill(prs, s, YELLOW)
            _add_text(s, Inches(1), Inches(2.7), Inches(11.3), Inches(1.6), sd.get("title", ""), 40, BLACK, bold=True)
            _add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.6), sd.get("subtitle", ""), 18, RGBColor(0x4A, 0x3A, 0x00))

        else:
            _bg_fill(prs, s, WHITE)
            _add_image_or_fallback(s, sd.get("_image"), Inches(8.3), Inches(0), Inches(5.03), prs.slide_height, PLACEHOLDER)
            tag = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(0.8), Inches(1.7), Inches(0.5))
            tag.fill.solid(); tag.fill.fore_color.rgb = YELLOW; _set_no_line(tag)
            _add_text(s, Inches(0.9), Inches(1.7), Inches(7), Inches(1.0), sd.get("title", ""), 26, BLACK, bold=True)
            _add_bullets(s, Inches(0.9), Inches(2.9), Inches(7), Inches(4.0), sd.get("bullets", []), 14, BLACK)


_RENDERERS = {
    "blue": _render_blue,
    "dark": _render_dark,
    "warm": _render_warm,
    "yellow": _render_yellow,
}


# ─────────────────────────────────────────────────────────
# Asosiy funksiya
# ─────────────────────────────────────────────────────────

def build_slide_deck(topic: str, num_slides: int, template: str) -> bytes:
    """Mavzu, slaydlar soni va shablon nomiga qarab to'liq PPTX faylni yasaydi."""
    with_images, _ = build_slide_deck_dual(topic, num_slides, template)
    return with_images


def build_slide_deck_dual(topic: str, num_slides: int, template: str) -> tuple:
    """Mavzu bo'yicha matn va rasmlarni BIR MARTA generatsiya qilib (Gemini/Pexels
    so'rovlari qayta yuborilmaydi), ikkita PPTX variant qaytaradi:
    (rasmli_bayt, rasmsiz_bayt)."""
    num_slides = max(MIN_SLIDES, min(MAX_SLIDES, num_slides))
    template = template if template in _RENDERERS else "blue"

    slides_data = generate_outline(topic, num_slides)

    # Har bir "content" slayd uchun Pexels'dan rasm topamiz
    for sd in slides_data:
        if sd.get("type", "content") == "content" and sd.get("image_keyword"):
            sd["_image"] = fetch_pexels_image(sd["image_keyword"])

    # ── Rasmli variant ────────────────────────────────────
    prs_with = Presentation()
    prs_with.slide_width = Inches(13.333)
    prs_with.slide_height = Inches(7.5)
    _RENDERERS[template](prs_with, slides_data)
    buf_with = io.BytesIO()
    prs_with.save(buf_with)
    buf_with.seek(0)

    # ── Rasmsiz variant - xuddi shu matn, lekin rasm o'rniga
    # brendga mos rangli blok (shablon dizayni buzilmaydi) ──
    slides_no_image = [dict(sd) for sd in slides_data]
    for sd in slides_no_image:
        sd["_image"] = None
    prs_without = Presentation()
    prs_without.slide_width = Inches(13.333)
    prs_without.slide_height = Inches(7.5)
    _RENDERERS[template](prs_without, slides_no_image)
    buf_without = io.BytesIO()
    prs_without.save(buf_without)
    buf_without.seek(0)

    return buf_with.read(), buf_without.read()
